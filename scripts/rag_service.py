"""Advanced RAG knowledge-base service used by the dashboard and analyzers.

MySQL stores auditable metadata and chunk text. LanceDB stores vectors. Cloud
calls are limited to embedding and reranking so the existing Ollama report
generation path remains unchanged.
"""

from __future__ import annotations

import csv
import hashlib
import io
import json
import math
import os
import re
import shutil
import time
import urllib.error
import urllib.request
import uuid
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional, Sequence, Tuple

import pymysql
from pymysql.cursors import DictCursor


DEFAULT_EMBEDDING_MODEL = "text-embedding-v4"
DEFAULT_RERANK_MODEL = "qwen3-rerank"
SUPPORTED_EXTENSIONS = {"txt", "md", "json", "jsonl", "csv", "pdf", "docx", "pptx", "xlsx"}
MAX_UPLOAD_BYTES = 30 * 1024 * 1024


def utc_text() -> str:
    return datetime.now().strftime("%Y-%m-%d %H:%M:%S")


def safe_filename(name: str) -> str:
    stem = re.sub(r"[^0-9A-Za-z._\-\u4e00-\u9fff]+", "_", Path(name or "document.txt").name)
    return stem[:160] or "document.txt"


def mysql_connect(conf: Dict[str, Any], autocommit: bool = False):
    return pymysql.connect(
        host=conf["host"],
        port=int(conf["port"]),
        user=conf["user"],
        password=conf["password"],
        database=conf["database"],
        charset="utf8mb4",
        cursorclass=DictCursor,
        autocommit=autocommit,
    )


def load_api_config(path: Path) -> Dict[str, Any]:
    payload: Dict[str, Any] = {}
    if path.exists():
        payload = json.loads(path.read_text(encoding="utf-8-sig"))
    api_key = str(os.environ.get("DASHSCOPE_API_KEY") or payload.get("api_key") or "").strip()
    base_url = str(
        os.environ.get("DASHSCOPE_BASE_URL")
        or payload.get("base_url")
        or "https://dashscope.aliyuncs.com/compatible-mode/v1"
    ).rstrip("/")
    rerank_url = str(
        os.environ.get("DASHSCOPE_RERANK_URL")
        or payload.get("rerank_url")
        or base_url.replace("/compatible-mode/v1", "/compatible-api/v1/reranks")
    ).rstrip("/")
    return {
        "api_key": api_key,
        "base_url": base_url,
        "rerank_url": rerank_url,
        "embedding_model": str(payload.get("embedding_model") or DEFAULT_EMBEDDING_MODEL),
        "rerank_model": str(payload.get("rerank_model") or DEFAULT_RERANK_MODEL),
        "embedding_dimensions": int(payload.get("embedding_dimensions") or 1024),
        "timeout_seconds": int(payload.get("timeout_seconds") or 90),
    }


def _post_json(url: str, api_key: str, payload: Dict[str, Any], timeout: int) -> Dict[str, Any]:
    if not api_key:
        raise RuntimeError("DashScope API Key 未配置")
    req = urllib.request.Request(
        url,
        data=json.dumps(payload, ensure_ascii=False).encode("utf-8"),
        headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
        method="POST",
    )
    try:
        with urllib.request.urlopen(req, timeout=timeout) as resp:
            return json.loads(resp.read().decode("utf-8"))
    except urllib.error.HTTPError as exc:
        detail = exc.read().decode("utf-8", errors="replace")[:500]
        raise RuntimeError(f"DashScope HTTP {exc.code}: {detail}") from exc
    except urllib.error.URLError as exc:
        raise RuntimeError(f"DashScope 连接失败: {exc.reason}") from exc


def embed_texts(texts: Sequence[str], config: Dict[str, Any]) -> List[List[float]]:
    clean = [str(x or "")[:8192] for x in texts]
    result: List[List[float]] = []
    for offset in range(0, len(clean), 10):
        batch = clean[offset : offset + 10]
        data = _post_json(
            f"{config['base_url']}/embeddings",
            config["api_key"],
            {"model": config["embedding_model"], "input": batch, "dimensions": config["embedding_dimensions"]},
            config["timeout_seconds"],
        )
        rows = sorted(data.get("data") or [], key=lambda x: int(x.get("index", 0)))
        if len(rows) != len(batch):
            raise RuntimeError("向量接口返回数量与输入不一致")
        result.extend([[float(v) for v in row["embedding"]] for row in rows])
    return result


def rerank(query: str, documents: Sequence[str], config: Dict[str, Any], top_n: int) -> List[Dict[str, Any]]:
    if not documents:
        return []
    payload = {
        "model": config["rerank_model"],
        "query": query[:4096],
        "documents": [str(x or "")[:8192] for x in documents],
        "top_n": min(max(1, top_n), len(documents)),
        "return_documents": False,
    }
    data = _post_json(config["rerank_url"], config["api_key"], payload, config["timeout_seconds"])
    return list((data.get("output") or {}).get("results") or data.get("results") or [])


def ensure_schema(conn: Any) -> None:
    statements = [
        """
        CREATE TABLE IF NOT EXISTS rag_knowledge_bases (
          id BIGINT PRIMARY KEY AUTO_INCREMENT,
          name VARCHAR(160) NOT NULL,
          description TEXT NULL,
          embedding_model VARCHAR(80) NOT NULL DEFAULT 'text-embedding-v4',
          rerank_model VARCHAR(80) NOT NULL DEFAULT 'qwen3-rerank',
          chunk_method VARCHAR(32) NOT NULL DEFAULT 'semantic',
          chunk_size INT NOT NULL DEFAULT 900,
          chunk_overlap INT NOT NULL DEFAULT 120,
          vector_top_k INT NOT NULL DEFAULT 20,
          keyword_top_k INT NOT NULL DEFAULT 20,
          final_top_k INT NOT NULL DEFAULT 5,
          score_threshold DECIMAL(8,6) NOT NULL DEFAULT 0.350000,
          enabled TINYINT(1) NOT NULL DEFAULT 1,
          scope_mode VARCHAR(20) NOT NULL DEFAULT 'all',
          created_by VARCHAR(80) NOT NULL DEFAULT 'system',
          created_at DATETIME NOT NULL DEFAULT CURRENT_TIMESTAMP,
          updated_at DATETIME NOT NULL DEFAULT CURRENT_TIMESTAMP ON UPDATE CURRENT_TIMESTAMP,
          UNIQUE KEY uk_rag_kb_name(name)
        ) ENGINE=InnoDB DEFAULT CHARSET=utf8mb4
        """,
        """
        CREATE TABLE IF NOT EXISTS rag_documents (
          id BIGINT PRIMARY KEY AUTO_INCREMENT,
          kb_id BIGINT NOT NULL,
          name VARCHAR(255) NOT NULL,
          source_type VARCHAR(32) NOT NULL DEFAULT 'upload',
          source_uri VARCHAR(1000) NULL,
          file_path VARCHAR(1000) NULL,
          mime_type VARCHAR(120) NULL,
          file_size BIGINT NOT NULL DEFAULT 0,
          checksum CHAR(64) NULL,
          status VARCHAR(24) NOT NULL DEFAULT 'pending',
          progress INT NOT NULL DEFAULT 0,
          chunk_count INT NOT NULL DEFAULT 0,
          char_count INT NOT NULL DEFAULT 0,
          error_message TEXT NULL,
          enabled TINYINT(1) NOT NULL DEFAULT 1,
          created_by VARCHAR(80) NOT NULL DEFAULT 'system',
          created_at DATETIME NOT NULL DEFAULT CURRENT_TIMESTAMP,
          updated_at DATETIME NOT NULL DEFAULT CURRENT_TIMESTAMP ON UPDATE CURRENT_TIMESTAMP,
          KEY idx_rag_document_kb(kb_id),
          KEY idx_rag_document_status(status)
        ) ENGINE=InnoDB DEFAULT CHARSET=utf8mb4
        """,
        """
        CREATE TABLE IF NOT EXISTS rag_chunks (
          id BIGINT PRIMARY KEY AUTO_INCREMENT,
          kb_id BIGINT NOT NULL,
          document_id BIGINT NOT NULL,
          chunk_index INT NOT NULL,
          title_path VARCHAR(600) NULL,
          content MEDIUMTEXT NOT NULL,
          token_count INT NOT NULL DEFAULT 0,
          vector_key VARCHAR(80) NOT NULL,
          enabled TINYINT(1) NOT NULL DEFAULT 1,
          retrieval_count BIGINT NOT NULL DEFAULT 0,
          last_retrieved_at DATETIME NULL,
          created_at DATETIME NOT NULL DEFAULT CURRENT_TIMESTAMP,
          updated_at DATETIME NOT NULL DEFAULT CURRENT_TIMESTAMP ON UPDATE CURRENT_TIMESTAMP,
          UNIQUE KEY uk_rag_chunk_doc_index(document_id, chunk_index),
          KEY idx_rag_chunk_kb(kb_id),
          FULLTEXT KEY ft_rag_chunk_content(title_path, content) WITH PARSER ngram
        ) ENGINE=InnoDB DEFAULT CHARSET=utf8mb4
        """,
        """
        CREATE TABLE IF NOT EXISTS rag_retrieval_tests (
          id BIGINT PRIMARY KEY AUTO_INCREMENT,
          kb_id BIGINT NOT NULL,
          query_text TEXT NOT NULL,
          duration_ms INT NOT NULL DEFAULT 0,
          result_count INT NOT NULL DEFAULT 0,
          results_json LONGTEXT NULL,
          created_by VARCHAR(80) NOT NULL DEFAULT 'system',
          created_at DATETIME NOT NULL DEFAULT CURRENT_TIMESTAMP,
          KEY idx_rag_test_kb_time(kb_id, created_at)
        ) ENGINE=InnoDB DEFAULT CHARSET=utf8mb4
        """,
        """
        CREATE TABLE IF NOT EXISTS rag_eval_cases (
          id BIGINT PRIMARY KEY AUTO_INCREMENT,
          kb_id BIGINT NOT NULL,
          question TEXT NOT NULL,
          expected_keywords TEXT NULL,
          expected_document VARCHAR(255) NULL,
          expected_no_match TINYINT(1) NOT NULL DEFAULT 0,
          expected_top_k INT NOT NULL DEFAULT 3,
          query_mode VARCHAR(20) NOT NULL DEFAULT 'incident',
          enabled TINYINT(1) NOT NULL DEFAULT 1,
          created_at DATETIME NOT NULL DEFAULT CURRENT_TIMESTAMP,
          KEY idx_rag_eval_kb(kb_id)
        ) ENGINE=InnoDB DEFAULT CHARSET=utf8mb4
        """,
        """
        CREATE TABLE IF NOT EXISTS rag_eval_runs (
          id BIGINT PRIMARY KEY AUTO_INCREMENT,
          kb_id BIGINT NOT NULL,
          total_cases INT NOT NULL DEFAULT 0,
          passed_cases INT NOT NULL DEFAULT 0,
          pass_rate DECIMAL(8,6) NOT NULL DEFAULT 0,
          average_duration_ms INT NOT NULL DEFAULT 0,
          recall_at_k DECIMAL(8,6) NOT NULL DEFAULT 0,
          mean_reciprocal_rank DECIMAL(8,6) NOT NULL DEFAULT 0,
          no_match_accuracy DECIMAL(8,6) NOT NULL DEFAULT 0,
          false_positive_rate DECIMAL(8,6) NOT NULL DEFAULT 0,
          config_json LONGTEXT NULL,
          results_json LONGTEXT NULL,
          created_by VARCHAR(80) NOT NULL DEFAULT 'system',
          created_at DATETIME NOT NULL DEFAULT CURRENT_TIMESTAMP,
          KEY idx_rag_eval_run_kb_time(kb_id, created_at)
        ) ENGINE=InnoDB DEFAULT CHARSET=utf8mb4
        """,
    ]
    with conn.cursor() as cur:
        for statement in statements:
            try:
                cur.execute(statement)
            except Exception as exc:
                # Some MySQL builds do not provide the ngram parser. Keep a
                # normal FULLTEXT index so the rest of the service still works.
                if "ngram" in statement.lower():
                    cur.execute(statement.replace(" WITH PARSER ngram", ""))
                else:
                    raise exc
        migrations = [
            "ALTER TABLE rag_eval_cases ADD COLUMN expected_no_match TINYINT(1) NOT NULL DEFAULT 0",
            "ALTER TABLE rag_eval_cases ADD COLUMN expected_top_k INT NOT NULL DEFAULT 3",
            "ALTER TABLE rag_eval_cases ADD COLUMN query_mode VARCHAR(20) NOT NULL DEFAULT 'incident'",
            "ALTER TABLE rag_eval_runs ADD COLUMN recall_at_k DECIMAL(8,6) NOT NULL DEFAULT 0",
            "ALTER TABLE rag_eval_runs ADD COLUMN mean_reciprocal_rank DECIMAL(8,6) NOT NULL DEFAULT 0",
            "ALTER TABLE rag_eval_runs ADD COLUMN no_match_accuracy DECIMAL(8,6) NOT NULL DEFAULT 0",
            "ALTER TABLE rag_eval_runs ADD COLUMN false_positive_rate DECIMAL(8,6) NOT NULL DEFAULT 0",
            "ALTER TABLE rag_eval_runs ADD COLUMN config_json LONGTEXT NULL",
        ]
        for statement in migrations:
            try:
                cur.execute(statement)
            except Exception as exc:
                if getattr(exc, "args", [None])[0] != 1060:
                    raise
        cur.execute("UPDATE rag_knowledge_bases SET score_threshold=0.350000 WHERE score_threshold < 0.350000")
    conn.commit()


def ensure_default_kb(conn: Any) -> int:
    with conn.cursor() as cur:
        cur.execute("SELECT id FROM rag_knowledge_bases ORDER BY id LIMIT 1")
        row = cur.fetchone()
        if row:
            return int(row["id"])
        cur.execute(
            """INSERT INTO rag_knowledge_bases
            (name, description, created_by) VALUES (%s, %s, %s)""",
            ("Web 攻击研判知识库", "用于 Web 攻击证据解释、风险判断与处置建议。", "migration"),
        )
        kb_id = int(cur.lastrowid)
    conn.commit()
    return kb_id


def parse_document(path: Path) -> str:
    ext = path.suffix.lower().lstrip(".")
    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"不支持的文件格式: .{ext}")
    if ext in {"txt", "md", "json", "jsonl", "csv"}:
        raw = path.read_bytes()
        text = raw.decode("utf-8-sig", errors="replace")
        if ext == "json":
            try:
                text = json.dumps(json.loads(text), ensure_ascii=False, indent=2)
            except Exception:
                pass
        return text
    if ext == "pdf":
        import fitz

        with fitz.open(path) as doc:
            return "\n\n".join(page.get_text("text") for page in doc)
    if ext == "docx":
        from docx import Document

        doc = Document(path)
        blocks = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                blocks.append(" | ".join(cell.text.strip() for cell in row.cells))
        return "\n".join(blocks)
    if ext == "pptx":
        from pptx import Presentation

        prs = Presentation(path)
        rows: List[str] = []
        for idx, slide in enumerate(prs.slides, 1):
            rows.append(f"# 幻灯片 {idx}")
            rows.extend(shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip())
        return "\n".join(rows)
    if ext == "xlsx":
        from openpyxl import load_workbook

        wb = load_workbook(path, read_only=True, data_only=True)
        rows = []
        for ws in wb.worksheets:
            rows.append(f"# 工作表 {ws.title}")
            for row in ws.iter_rows(values_only=True):
                values = [str(v).strip() for v in row if v is not None and str(v).strip()]
                if values:
                    rows.append(" | ".join(values))
        return "\n".join(rows)
    return ""


def _split_long_text(text: str, size: int, overlap: int) -> List[str]:
    text = re.sub(r"[ \t]+", " ", text).strip()
    if not text:
        return []
    chunks: List[str] = []
    start = 0
    while start < len(text):
        end = min(len(text), start + size)
        if end < len(text):
            candidates = [text.rfind(mark, start + size // 2, end) for mark in ("\n\n", "。", "；", "\n", ". ")]
            boundary = max(candidates)
            if boundary > start:
                end = boundary + 1
        piece = text[start:end].strip()
        if piece:
            chunks.append(piece)
        if end >= len(text):
            break
        start = max(start + 1, end - overlap)
    return chunks


def chunk_document(text: str, method: str = "semantic", size: int = 900, overlap: int = 120) -> List[Dict[str, Any]]:
    size = min(max(int(size), 200), 4000)
    overlap = min(max(int(overlap), 0), size // 2)
    text = text.replace("\r\n", "\n").replace("\r", "\n").strip()
    if not text:
        return []
    if method == "fixed":
        parts = _split_long_text(text, size, overlap)
        return [{"title_path": "", "content": part} for part in parts]
    heading = "正文"
    sections: List[Tuple[str, str]] = []
    buffer: List[str] = []
    for line in text.split("\n"):
        clean = line.strip()
        is_heading = bool(re.match(r"^#{1,6}\s+", clean) or re.match(r"^[一二三四五六七八九十0-9]+[、.．]\s*\S+", clean))
        if is_heading and buffer:
            sections.append((heading, "\n".join(buffer).strip()))
            buffer = []
        if is_heading:
            heading = re.sub(r"^#{1,6}\s+", "", clean)[:300]
        else:
            buffer.append(line)
    if buffer:
        sections.append((heading, "\n".join(buffer).strip()))
    chunks: List[Dict[str, Any]] = []
    for title, body in sections or [("正文", text)]:
        for piece in _split_long_text(body, size, overlap):
            chunks.append({"title_path": title, "content": piece})
    return chunks


def _lance_table(data_dir: Path, vector_size: int, create: bool = False):
    import lancedb
    import pyarrow as pa

    db = lancedb.connect(str(data_dir / "lancedb"))
    if "rag_chunks" in db.table_names():
        return db.open_table("rag_chunks")
    if not create:
        return None
    schema = pa.schema(
        [
            pa.field("vector_key", pa.string()),
            pa.field("kb_id", pa.int64()),
            pa.field("document_id", pa.int64()),
            pa.field("chunk_id", pa.int64()),
            pa.field("title", pa.string()),
            pa.field("content", pa.string()),
            pa.field("vector", pa.list_(pa.float32(), vector_size)),
        ]
    )
    return db.create_table("rag_chunks", schema=schema)


def _delete_vectors(data_dir: Path, where: str) -> None:
    table = _lance_table(data_dir, 1024, create=False)
    if table is not None:
        try:
            table.delete(where)
        except Exception:
            pass


def list_kbs(conn: Any, q: str = "", include_disabled: bool = True) -> List[Dict[str, Any]]:
    where = []
    params: List[Any] = []
    if q:
        where.append("(k.name LIKE %s OR k.description LIKE %s)")
        params.extend([f"%{q}%", f"%{q}%"])
    if not include_disabled:
        where.append("k.enabled=1")
    clause = f"WHERE {' AND '.join(where)}" if where else ""
    with conn.cursor() as cur:
        cur.execute(
            f"""SELECT k.*,
            (SELECT COUNT(*) FROM rag_documents d WHERE d.kb_id=k.id) AS document_count,
            (SELECT COUNT(*) FROM rag_chunks c WHERE c.kb_id=k.id AND c.enabled=1) AS chunk_count
            FROM rag_knowledge_bases k {clause} ORDER BY k.updated_at DESC, k.id DESC""",
            tuple(params),
        )
        return list(cur.fetchall())


def get_kb(conn: Any, kb_id: int) -> Optional[Dict[str, Any]]:
    rows = [row for row in list_kbs(conn) if int(row["id"]) == int(kb_id)]
    return rows[0] if rows else None


def save_kb(conn: Any, body: Dict[str, Any], username: str, kb_id: Optional[int] = None) -> int:
    name = str(body.get("name") or "").strip()
    if not name:
        raise ValueError("知识库名称不能为空")
    values = (
        name,
        str(body.get("description") or "").strip(),
        str(body.get("embedding_model") or DEFAULT_EMBEDDING_MODEL),
        str(body.get("rerank_model") or DEFAULT_RERANK_MODEL),
        "fixed" if body.get("chunk_method") == "fixed" else "semantic",
        min(max(int(body.get("chunk_size") or 900), 200), 4000),
        min(max(int(body.get("chunk_overlap") or 120), 0), 1000),
        min(max(int(body.get("vector_top_k") or 20), 1), 100),
        min(max(int(body.get("keyword_top_k") or 20), 1), 100),
        min(max(int(body.get("final_top_k") or 5), 1), 20),
        max(0.0, min(float(body.get("score_threshold") or 0.35), 1.0)),
        1 if body.get("enabled", True) else 0,
    )
    with conn.cursor() as cur:
        if kb_id:
            cur.execute(
                """UPDATE rag_knowledge_bases SET name=%s,description=%s,embedding_model=%s,rerank_model=%s,
                chunk_method=%s,chunk_size=%s,chunk_overlap=%s,vector_top_k=%s,keyword_top_k=%s,
                final_top_k=%s,score_threshold=%s,enabled=%s WHERE id=%s""",
                values + (int(kb_id),),
            )
            if cur.rowcount == 0:
                raise KeyError("知识库不存在")
            result = int(kb_id)
        else:
            cur.execute(
                """INSERT INTO rag_knowledge_bases
                (name,description,embedding_model,rerank_model,chunk_method,chunk_size,chunk_overlap,
                vector_top_k,keyword_top_k,final_top_k,score_threshold,enabled,created_by)
                VALUES (%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s)""",
                values + (username,),
            )
            result = int(cur.lastrowid)
    conn.commit()
    return result


def delete_kb(conn: Any, data_dir: Path, kb_id: int) -> bool:
    with conn.cursor() as cur:
        cur.execute("SELECT file_path FROM rag_documents WHERE kb_id=%s", (kb_id,))
        files = [str(row.get("file_path") or "") for row in cur.fetchall()]
        cur.execute("DELETE FROM rag_retrieval_tests WHERE kb_id=%s", (kb_id,))
        cur.execute("DELETE FROM rag_eval_cases WHERE kb_id=%s", (kb_id,))
        cur.execute("DELETE FROM rag_eval_runs WHERE kb_id=%s", (kb_id,))
        cur.execute("DELETE FROM rag_chunks WHERE kb_id=%s", (kb_id,))
        cur.execute("DELETE FROM rag_documents WHERE kb_id=%s", (kb_id,))
        cur.execute("DELETE FROM rag_knowledge_bases WHERE id=%s", (kb_id,))
        changed = cur.rowcount > 0
    conn.commit()
    _delete_vectors(data_dir, f"kb_id = {int(kb_id)}")
    for file_path in files:
        try:
            Path(file_path).unlink(missing_ok=True)
        except Exception:
            pass
    return changed


def list_documents(conn: Any, kb_id: int) -> List[Dict[str, Any]]:
    with conn.cursor() as cur:
        cur.execute("SELECT * FROM rag_documents WHERE kb_id=%s ORDER BY updated_at DESC,id DESC", (kb_id,))
        return list(cur.fetchall())


def list_chunks(conn: Any, kb_id: int, document_id: Optional[int] = None) -> List[Dict[str, Any]]:
    sql = """SELECT c.*,d.name AS document_name FROM rag_chunks c
             JOIN rag_documents d ON d.id=c.document_id WHERE c.kb_id=%s"""
    params: List[Any] = [kb_id]
    if document_id:
        sql += " AND c.document_id=%s"
        params.append(document_id)
    sql += " ORDER BY c.document_id,c.chunk_index"
    with conn.cursor() as cur:
        cur.execute(sql, tuple(params))
        return list(cur.fetchall())


def ingest_file(
    conn: Any,
    data_dir: Path,
    api_config: Dict[str, Any],
    kb_id: int,
    source_path: Path,
    original_name: str,
    username: str,
    source_type: str = "upload",
    source_uri: str = "",
) -> Dict[str, Any]:
    kb = get_kb(conn, kb_id)
    if not kb:
        raise KeyError("知识库不存在")
    size = source_path.stat().st_size
    if size > MAX_UPLOAD_BYTES:
        raise ValueError("文件不能超过 30MB")
    extension = source_path.suffix.lower().lstrip(".")
    if extension not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"不支持 .{extension} 文件")
    raw = source_path.read_bytes()
    checksum = hashlib.sha256(raw).hexdigest()
    with conn.cursor() as cur:
        cur.execute(
            """SELECT id,name,status,chunk_count,char_count FROM rag_documents
            WHERE kb_id=%s AND checksum=%s AND status='ready' ORDER BY id DESC LIMIT 1""",
            (kb_id, checksum),
        )
        existing = cur.fetchone()
    if existing:
        return {**existing, "skipped": True, "reason": "相同内容已经入库"}
    target_dir = data_dir / "uploads" / str(kb_id)
    target_dir.mkdir(parents=True, exist_ok=True)
    target_path = target_dir / f"{uuid.uuid4().hex}_{safe_filename(original_name)}"
    shutil.copy2(source_path, target_path)
    with conn.cursor() as cur:
        cur.execute(
            """INSERT INTO rag_documents
            (kb_id,name,source_type,source_uri,file_path,mime_type,file_size,checksum,status,progress,created_by)
            VALUES (%s,%s,%s,%s,%s,%s,%s,%s,'processing',10,%s)""",
            (kb_id, original_name[:255], source_type, source_uri[:1000], str(target_path), extension, size, checksum, username),
        )
        document_id = int(cur.lastrowid)
    conn.commit()
    try:
        text = parse_document(target_path)
        chunks = chunk_document(text, str(kb["chunk_method"]), int(kb["chunk_size"]), int(kb["chunk_overlap"]))
        if not chunks:
            raise ValueError("文档未解析出可用文本")
        vectors: List[List[float]] = []
        for offset in range(0, len(chunks), 10):
            vectors.extend(embed_texts([row["content"] for row in chunks[offset : offset + 10]], api_config))
            with conn.cursor() as cur:
                cur.execute(
                    "UPDATE rag_documents SET progress=%s WHERE id=%s",
                    (min(85, 20 + int((offset + 10) / max(1, len(chunks)) * 65)), document_id),
                )
            conn.commit()
        vector_rows = []
        with conn.cursor() as cur:
            for index, (chunk, vector) in enumerate(zip(chunks, vectors)):
                vector_key = uuid.uuid4().hex
                cur.execute(
                    """INSERT INTO rag_chunks
                    (kb_id,document_id,chunk_index,title_path,content,token_count,vector_key)
                    VALUES (%s,%s,%s,%s,%s,%s,%s)""",
                    (kb_id, document_id, index, chunk["title_path"], chunk["content"], max(1, len(chunk["content"]) // 2), vector_key),
                )
                chunk_id = int(cur.lastrowid)
                vector_rows.append(
                    {
                        "vector_key": vector_key,
                        "kb_id": int(kb_id),
                        "document_id": document_id,
                        "chunk_id": chunk_id,
                        "title": str(chunk["title_path"] or original_name),
                        "content": chunk["content"],
                        "vector": vector,
                    }
                )
            cur.execute(
                """UPDATE rag_documents SET status='ready',progress=100,chunk_count=%s,char_count=%s,error_message=NULL
                WHERE id=%s""",
                (len(chunks), len(text), document_id),
            )
        conn.commit()
        table = _lance_table(data_dir, int(api_config["embedding_dimensions"]), create=True)
        table.add(vector_rows)
        return {"id": document_id, "name": original_name, "status": "ready", "chunk_count": len(chunks), "char_count": len(text)}
    except Exception as exc:
        with conn.cursor() as cur:
            cur.execute(
                "UPDATE rag_documents SET status='failed',progress=100,error_message=%s WHERE id=%s",
                (str(exc)[:2000], document_id),
            )
        conn.commit()
        raise


def delete_document(conn: Any, data_dir: Path, document_id: int) -> bool:
    with conn.cursor() as cur:
        cur.execute("SELECT file_path FROM rag_documents WHERE id=%s", (document_id,))
        row = cur.fetchone()
        if not row:
            return False
        cur.execute("DELETE FROM rag_chunks WHERE document_id=%s", (document_id,))
        cur.execute("DELETE FROM rag_documents WHERE id=%s", (document_id,))
    conn.commit()
    _delete_vectors(data_dir, f"document_id = {int(document_id)}")
    try:
        Path(str(row.get("file_path") or "")).unlink(missing_ok=True)
    except Exception:
        pass
    return True


def update_chunk(conn: Any, data_dir: Path, api_config: Dict[str, Any], chunk_id: int, content: str, enabled: bool) -> None:
    content = str(content or "").strip()
    if not content:
        raise ValueError("切片内容不能为空")
    with conn.cursor() as cur:
        cur.execute("SELECT * FROM rag_chunks WHERE id=%s", (chunk_id,))
        row = cur.fetchone()
        if not row:
            raise KeyError("切片不存在")
        cur.execute(
            "UPDATE rag_chunks SET content=%s,enabled=%s,token_count=%s WHERE id=%s",
            (content, 1 if enabled else 0, max(1, len(content) // 2), chunk_id),
        )
    conn.commit()
    vector = embed_texts([content], api_config)[0]
    table = _lance_table(data_dir, len(vector), create=True)
    table.delete(f"chunk_id = {int(chunk_id)}")
    table.add(
        [
            {
                "vector_key": str(row["vector_key"]),
                "kb_id": int(row["kb_id"]),
                "document_id": int(row["document_id"]),
                "chunk_id": int(chunk_id),
                "title": str(row.get("title_path") or ""),
                "content": content,
                "vector": vector,
            }
        ]
    )


def _tokens(text: str) -> List[str]:
    try:
        import jieba

        values = jieba.lcut(str(text or "").lower(), cut_all=False)
    except Exception:
        values = re.findall(r"[a-z0-9_\-\.]+|[\u4e00-\u9fff]", str(text or "").lower())
    return [x.strip() for x in values if x.strip() and not x.isspace()]


def _bm25_scores(query: str, rows: Sequence[Dict[str, Any]]) -> Dict[int, float]:
    query_tokens = _tokens(query)
    docs = [_tokens(f"{row.get('title_path','')} {row.get('content','')}") for row in rows]
    if not query_tokens or not docs:
        return {}
    avgdl = sum(len(doc) for doc in docs) / max(1, len(docs))
    df: Dict[str, int] = {}
    for doc in docs:
        for token in set(doc):
            df[token] = df.get(token, 0) + 1
    scores: Dict[int, float] = {}
    for row, doc in zip(rows, docs):
        freq: Dict[str, int] = {}
        for token in doc:
            freq[token] = freq.get(token, 0) + 1
        score = 0.0
        for token in query_tokens:
            tf = freq.get(token, 0)
            if not tf:
                continue
            idf = math.log(1 + (len(docs) - df.get(token, 0) + 0.5) / (df.get(token, 0) + 0.5))
            score += idf * (tf * 2.2) / (tf + 1.2 * (1 - 0.75 + 0.75 * len(doc) / max(1.0, avgdl)))
        if score > 0:
            scores[int(row["id"])] = score
    return scores


def expand_security_query(query: str) -> str:
    """Add canonical security terms for observable evidence descriptions."""
    raw = str(query or "").strip()
    lowered = raw.lower()
    additions: List[str] = []
    rules = [
        (("nmap" in lowered) or ("syn" in lowered and "端口" in raw), "PortScan port scan 端口扫描 reconnaissance"),
        (("404" in lowered and ("路径" in raw or "目录" in raw)) or "目录枚举" in raw, "directory scan path enumeration 目录扫描"),
        (("or 1=1" in lowered) or ("单引号" in raw and "注释符" in raw), "SQL Injection SQL 注入 parameterized query"),
        (("<script" in lowered) or "onerror" in lowered or ("html" in lowered and "用户输入" in raw), "XSS cross site scripting output encoding CSP"),
        (("爆破" in raw) or ("失败登录" in raw and "短时间" in raw), "brute force credential attack 暴力破解"),
        (("../" in raw) or "路径穿越" in raw, "path traversal directory traversal 路径遍历"),
        (("账号锁定" in raw) or ("失败计数" in raw) or "account lockout" in lowered, "account lockout authentication throttling failed login counter 账号锁定"),
    ]
    for matched, terms in rules:
        if matched:
            additions.append(terms)
    return f"{raw}\n安全领域标准术语：{' '.join(additions)}" if additions else raw


def analyze_security_query(query: str) -> Dict[str, Any]:
    """Classify retrieval intent without treating retrieval as attack detection."""
    raw = str(query or "").strip()
    lowered = raw.lower()
    compact = re.sub(r"\s+", "", raw)
    tokens = _tokens(raw)
    alpha_numeric = re.findall(r"[a-z0-9\u4e00-\u9fff]", lowered)
    unique_ratio = len(set(alpha_numeric)) / max(1, len(alpha_numeric))
    ascii_blob = bool(re.fullmatch(r"[a-z0-9_-]+", lowered))
    gibberish = bool(
        len(compact) < 4
        or not alpha_numeric
        or (ascii_blob and len(compact) >= 16 and unique_ratio > 0.60)
    )
    security_patterns = [
        r"\b(?:sql|xss|csrf|ssrf|rce|xxe|owasp|mitre|nmap|payload|shell|security|authentication|authorization|credential|hardening|mitigation|firewall|waf)\b|cve[-_ ]?\d+|vulnerabilit|account\s+lockout|漏洞|注入|跨站|扫描|爆破|攻击|防护|认证|登录失败|密码|端口|目录|多个路径|路径穿越|账号锁定|失败计数|解锁策略|恶意|告警",
        r"union\s+select|or\s+['\"]?1['\"]?\s*=\s*['\"]?1|<script|onerror\s*=|\.\./|\$\{jndi:|/etc/passwd|cmd\.exe|powershell",
    ]
    incident_patterns = [
        r"\b(?:get|post|put|delete|patch|head|options)\s+/|\buri\s*=|\brequest(?:_body)?\s*=|\bstatus(?:_code)?\s*=|响应.*用户输入|response.*user\s+input",
        r"\b(?:\d{1,3}\.){3}\d{1,3}\b|\bport\s*[=:]?\s*\d{1,5}\b|端口\s*\d{1,5}",
        r"union\s+select|or\s+['\"]?1['\"]?\s*=\s*['\"]?1|<script|script\s*标签|onerror(?:\s*=)?|\.\./|\$\{jndi:|/etc/passwd|cmd\.exe|powershell",
        r"短时间|连续|大量|多次|多个端口|多个路径|失败\s*\d+\s*次|请求\s*\d+\s*次|syn\s+(?:probe|scan)|directory\s+(?:scan|enumeration)|brute\s*force",
    ]
    attack_indicator_patterns = [
        r"sql\s*注入|sql\s*injection|\b(?:xss|csrf|ssrf|rce|xxe|nmap)\b|cross[- ]site|port\s*scan|directory\s+(?:scan|enumeration)|brute\s*force|漏洞利用|端口扫描|目录扫描|暴力破解|恶意|攻击",
        r"union\s+select|or\s+['\"]?1['\"]?\s*=\s*['\"]?1|<script|script\s*标签|onerror(?:\s*=)?|\.\./|\$\{jndi:|/etc/passwd|cmd\.exe|powershell",
        r"(?:短时间|一分钟内|\d+\s*秒内).*(?:连续|大量|多次|端口|路径|失败)|(?:连续|大量|多次|多个端口|多个路径).*(?:短时间|一分钟内|\d+\s*秒内)|多个路径.*大量|连续\s*404|失败\s*[5-9]\d*\s*次|请求\s*[1-9]\d+\s*次|syn\s+(?:probe|scan)",
    ]
    security_hits = sum(1 for pattern in security_patterns if re.search(pattern, lowered, re.I))
    incident_hits = sum(1 for pattern in incident_patterns if re.search(pattern, lowered, re.I))
    attack_hits = sum(1 for pattern in attack_indicator_patterns if re.search(pattern, lowered, re.I))
    return {
        "gibberish": gibberish,
        "token_count": len(tokens),
        "security_signal_count": security_hits,
        "incident_evidence_count": incident_hits,
        "attack_indicator_count": attack_hits,
        "has_security_context": security_hits > 0 or incident_hits > 0,
        "has_incident_evidence": incident_hits > 0,
        "has_attack_indicator": attack_hits > 0,
    }


def retrieval_gate(query: str, requested_mode: str = "auto") -> Dict[str, Any]:
    profile = analyze_security_query(query)
    mode = requested_mode if requested_mode in {"auto", "incident", "knowledge"} else "auto"
    if mode == "auto":
        mode = "incident" if profile["has_incident_evidence"] else "knowledge"
    reason = ""
    if profile["gibberish"]:
        reason = "输入缺少可识别语义，已停止无效召回"
    elif mode == "incident" and not profile["has_incident_evidence"]:
        reason = "未发现请求、Payload、IP/端口或时间窗口行为证据"
    elif mode == "incident" and not profile["has_attack_indicator"]:
        reason = "事件中没有可复核的攻击特征，禁止仅凭语义相近生成攻击知识"
    elif mode == "knowledge" and not profile["has_security_context"]:
        reason = "问题与当前网络安全知识域无明确关联"
    return {"accepted": not reason, "mode": mode, "reason": reason, "profile": profile}


def diversify_candidates(rows: Sequence[Dict[str, Any]], per_document: int = 12, limit: int = 30) -> List[Dict[str, Any]]:
    selected: List[Dict[str, Any]] = []
    counts: Dict[int, int] = {}
    for row in rows:
        document_id = int(row.get("document_id") or 0)
        if counts.get(document_id, 0) >= per_document:
            continue
        selected.append(row)
        counts[document_id] = counts.get(document_id, 0) + 1
        if len(selected) >= limit:
            break
    return selected


def hybrid_search(
    conn: Any,
    data_dir: Path,
    api_config: Dict[str, Any],
    kb_id: int,
    query: str,
    username: str = "system",
    save_test: bool = False,
    query_mode: str = "auto",
) -> Dict[str, Any]:
    started = time.perf_counter()
    kb = get_kb(conn, kb_id)
    if not kb:
        raise KeyError("知识库不存在")
    query = str(query or "").strip()
    if not query:
        raise ValueError("测试文本不能为空")
    rows = list_chunks(conn, kb_id)
    rows = [row for row in rows if int(row.get("enabled") or 0) == 1]
    if not rows:
        return {"items": [], "duration_ms": 0, "query": query, "no_match": True, "no_match_reason": "知识库中没有启用的切片"}
    gate = retrieval_gate(query, query_mode)
    effective_threshold = max(0.35, float(kb["score_threshold"]))
    if not gate["accepted"]:
        duration_ms = int((time.perf_counter() - started) * 1000)
        result = {
            "items": [], "duration_ms": duration_ms, "query": query, "query_mode": gate["mode"],
            "no_match": True, "no_match_reason": gate["reason"], "effective_threshold": effective_threshold,
            "score_semantics": "重排匹配分，不代表攻击概率或准确率", "query_profile": gate["profile"],
        }
        if save_test:
            with conn.cursor() as cur:
                cur.execute(
                    """INSERT INTO rag_retrieval_tests
                    (kb_id,query_text,duration_ms,result_count,results_json,created_by) VALUES (%s,%s,%s,%s,%s,%s)""",
                    (kb_id, query, duration_ms, 0, json.dumps(result, ensure_ascii=False), username),
                )
            conn.commit()
        return result
    search_query = expand_security_query(query)
    vector_scores: Dict[int, float] = {}
    vector_distances: Dict[int, float] = {}
    table = _lance_table(data_dir, int(api_config["embedding_dimensions"]), create=False)
    if table is not None:
        vector = embed_texts([search_query], api_config)[0]
        result = (
            table.search(vector)
            .where(f"kb_id = {int(kb_id)}")
            .limit(min(int(kb["vector_top_k"]), len(rows)))
            .to_list()
        )
        for rank, item in enumerate(result, 1):
            vector_scores[int(item["chunk_id"])] = 1.0 / (60 + rank)
            if item.get("_distance") is not None:
                vector_distances[int(item["chunk_id"])] = float(item["_distance"])
    bm25 = _bm25_scores(search_query, rows)
    keyword_ranked = sorted(bm25, key=bm25.get, reverse=True)[: int(kb["keyword_top_k"])]
    fused: Dict[int, float] = dict(vector_scores)
    for rank, chunk_id in enumerate(keyword_ranked, 1):
        fused[chunk_id] = fused.get(chunk_id, 0.0) + 1.0 / (60 + rank)
    row_map = {int(row["id"]): row for row in rows}
    fused_rows = [row_map[cid] for cid in sorted(fused, key=fused.get, reverse=True) if cid in row_map]
    candidates = diversify_candidates(fused_rows, per_document=12, limit=30)
    ranked = rerank(search_query, [row["content"] for row in candidates], api_config, min(20, len(candidates)))
    items = []
    document_counts: Dict[int, int] = {}
    threshold = effective_threshold
    top_candidate_score = 0.0
    for item in ranked:
        index = int(item.get("index", -1))
        if index < 0 or index >= len(candidates):
            continue
        row = dict(candidates[index])
        document_id = int(row["document_id"])
        if document_counts.get(document_id, 0) >= 3:
            continue
        score = float(item.get("relevance_score") or item.get("score") or 0.0)
        top_candidate_score = max(top_candidate_score, score)
        if score < threshold:
            continue
        chunk_id = int(row["id"])
        lexical_score = float(bm25.get(chunk_id, 0.0))
        # Mid-confidence matches need either lexical support or a strong rerank
        # score. This prevents vector-nearest but unrelated text from leaking in.
        if score < 0.55 and lexical_score <= 0:
            continue
        items.append(
            {
                "chunk_id": int(row["id"]),
                "document_id": int(row["document_id"]),
                "document_name": row.get("document_name"),
                "chunk_index": int(row["chunk_index"]),
                "title_path": row.get("title_path") or "正文",
                "content": row.get("content") or "",
                "score": round(score, 6),
                "rrf_score": round(fused.get(chunk_id, 0.0), 6),
                "bm25_score": round(lexical_score, 6),
                "vector_distance": round(vector_distances[chunk_id], 6) if chunk_id in vector_distances else None,
            }
        )
        document_counts[document_id] = document_counts.get(document_id, 0) + 1
        if len(items) >= int(kb["final_top_k"]):
            break
    duration_ms = int((time.perf_counter() - started) * 1000)
    if items:
        ids = [item["chunk_id"] for item in items]
        placeholders = ",".join(["%s"] * len(ids))
        with conn.cursor() as cur:
            cur.execute(
                f"UPDATE rag_chunks SET retrieval_count=retrieval_count+1,last_retrieved_at=NOW() WHERE id IN ({placeholders})",
                tuple(ids),
            )
        conn.commit()
    no_match_reason = ""
    if not items:
        if top_candidate_score and top_candidate_score < threshold:
            no_match_reason = f"最高重排分 {top_candidate_score:.3f} 低于有效阈值 {threshold:.3f}"
        else:
            no_match_reason = "候选结果缺少足够的语义与关键词交叉证据"
    response = {
        "items": items, "duration_ms": duration_ms, "query": query, "query_mode": gate["mode"],
        "no_match": not items, "no_match_reason": no_match_reason,
        "effective_threshold": threshold, "top_candidate_score": round(top_candidate_score, 6),
        "score_semantics": "重排匹配分，不代表攻击概率或准确率", "query_profile": gate["profile"],
    }
    if save_test:
        with conn.cursor() as cur:
            cur.execute(
                """INSERT INTO rag_retrieval_tests
                (kb_id,query_text,duration_ms,result_count,results_json,created_by) VALUES (%s,%s,%s,%s,%s,%s)""",
                (kb_id, query, duration_ms, len(items), json.dumps(response, ensure_ascii=False), username),
            )
        conn.commit()
    return response


def list_test_history(conn: Any, kb_id: int, limit: int = 30) -> List[Dict[str, Any]]:
    with conn.cursor() as cur:
        cur.execute(
            """SELECT id,kb_id,query_text,duration_ms,result_count,created_by,created_at
            FROM rag_retrieval_tests WHERE kb_id=%s ORDER BY id DESC LIMIT %s""",
            (kb_id, min(max(1, limit), 100)),
        )
        return list(cur.fetchall())


def list_eval_cases(conn: Any, kb_id: int) -> List[Dict[str, Any]]:
    with conn.cursor() as cur:
        cur.execute(
            """SELECT id,kb_id,question,expected_keywords,expected_document,expected_no_match,
            expected_top_k,query_mode,enabled,created_at
            FROM rag_eval_cases WHERE kb_id=%s ORDER BY id""",
            (kb_id,),
        )
        return list(cur.fetchall())


def save_eval_case(conn: Any, kb_id: int, payload: Dict[str, Any], case_id: int = 0) -> Dict[str, Any]:
    question = str(payload.get("question") or "").strip()
    if not question:
        raise ValueError("测试问题不能为空")
    keywords = str(payload.get("expected_keywords") or "").strip()
    expected_document = str(payload.get("expected_document") or "").strip()
    expected_no_match = 1 if bool(payload.get("expected_no_match", False)) else 0
    expected_top_k = min(max(int(payload.get("expected_top_k") or 3), 1), 20)
    query_mode = str(payload.get("query_mode") or "incident").strip().lower()
    if query_mode not in {"incident", "knowledge", "auto"}:
        query_mode = "incident"
    enabled = 1 if bool(payload.get("enabled", True)) else 0
    with conn.cursor() as cur:
        if case_id:
            cur.execute(
                """UPDATE rag_eval_cases SET question=%s,expected_keywords=%s,
                expected_document=%s,expected_no_match=%s,expected_top_k=%s,query_mode=%s,
                enabled=%s WHERE id=%s AND kb_id=%s""",
                (question, keywords, expected_document, expected_no_match, expected_top_k, query_mode, enabled, case_id, kb_id),
            )
            if not cur.rowcount:
                raise KeyError("回归用例不存在")
        else:
            cur.execute(
                """INSERT INTO rag_eval_cases
                (kb_id,question,expected_keywords,expected_document,expected_no_match,expected_top_k,query_mode,enabled)
                VALUES (%s,%s,%s,%s,%s,%s,%s,%s)""",
                (kb_id, question, keywords, expected_document, expected_no_match, expected_top_k, query_mode, enabled),
            )
            case_id = int(cur.lastrowid)
    conn.commit()
    return next(row for row in list_eval_cases(conn, kb_id) if int(row["id"]) == int(case_id))


def delete_eval_case(conn: Any, kb_id: int, case_id: int) -> bool:
    with conn.cursor() as cur:
        cur.execute("DELETE FROM rag_eval_cases WHERE id=%s AND kb_id=%s", (case_id, kb_id))
        changed = bool(cur.rowcount)
    conn.commit()
    return changed


def evaluate_retrieval_items(
    items: Sequence[Dict[str, Any]], expected_keywords: str = "", expected_document: str = "",
    expected_no_match: bool = False, expected_top_k: int = 3,
) -> Dict[str, Any]:
    keywords = [part.strip().lower() for part in re.split(r"[,，;；\n]+", expected_keywords or "") if part.strip()]
    if expected_no_match:
        passed = not items
        return {
            "passed": passed, "matched_keywords": [], "relevant_rank": 0, "reciprocal_rank": 0.0,
            "reason": "正确拒绝无关查询" if passed else f"期望无结果，但错误召回 {len(items)} 个切片",
        }
    expected_top_k = min(max(int(expected_top_k or 3), 1), 20)
    needle = expected_document.strip().lower()
    relevant_rank = 0
    matched_keywords: List[str] = []
    for rank, item in enumerate(items, 1):
        combined = f"{item.get('document_name') or ''}\n{item.get('title_path') or ''}\n{item.get('content') or ''}".lower()
        item_keywords = [keyword for keyword in keywords if keyword in combined]
        document_ok = not needle or needle in str(item.get("document_name") or "").lower()
        keyword_ok = not keywords or bool(item_keywords)
        if document_ok and keyword_ok:
            relevant_rank = rank
            matched_keywords = item_keywords
            break
    passed = bool(relevant_rank and relevant_rank <= expected_top_k)
    reasons = []
    if not items:
        reasons.append("未召回任何切片")
    elif not relevant_rank:
        if keywords:
            reasons.append("没有单个切片命中期望关键词")
        if needle:
            reasons.append("未召回同时满足条件的期望来源文档")
    elif relevant_rank > expected_top_k:
        reasons.append(f"正确证据位于 TOP{relevant_rank}，低于 TOP{expected_top_k} 要求")
    return {
        "passed": passed,
        "matched_keywords": matched_keywords,
        "relevant_rank": relevant_rank,
        "reciprocal_rank": round(1.0 / relevant_rank, 6) if relevant_rank else 0.0,
        "reason": "；".join(reasons) if reasons else f"正确证据位于 TOP{relevant_rank}",
    }


def run_eval_suite(
    conn: Any,
    data_dir: Path,
    api_config: Dict[str, Any],
    kb_id: int,
    username: str = "system",
) -> Dict[str, Any]:
    cases = [row for row in list_eval_cases(conn, kb_id) if int(row.get("enabled") or 0) == 1]
    if not cases:
        raise ValueError("请先添加至少一条已启用的回归用例")
    results: List[Dict[str, Any]] = []
    for case in cases:
        search_result = hybrid_search(
            conn, data_dir, api_config, kb_id, str(case["question"]), username=username, save_test=False,
            query_mode=str(case.get("query_mode") or "incident"),
        )
        evaluation = evaluate_retrieval_items(
            search_result.get("items") or [],
            str(case.get("expected_keywords") or ""),
            str(case.get("expected_document") or ""),
            bool(case.get("expected_no_match")),
            int(case.get("expected_top_k") or 3),
        )
        results.append(
            {
                "case_id": int(case["id"]),
                "question": case["question"],
                "expected_keywords": case.get("expected_keywords") or "",
                "expected_document": case.get("expected_document") or "",
                "expected_no_match": bool(case.get("expected_no_match")),
                "expected_top_k": int(case.get("expected_top_k") or 3),
                "query_mode": case.get("query_mode") or "incident",
                "duration_ms": int(search_result.get("duration_ms") or 0),
                "no_match": bool(search_result.get("no_match")),
                "no_match_reason": search_result.get("no_match_reason") or "",
                "top_document": (search_result.get("items") or [{}])[0].get("document_name") if search_result.get("items") else "",
                "top_score": (search_result.get("items") or [{}])[0].get("score") if search_result.get("items") else 0,
                **evaluation,
            }
        )
    passed = sum(1 for row in results if row["passed"])
    total = len(results)
    average_ms = int(sum(row["duration_ms"] for row in results) / total) if total else 0
    pass_rate = passed / total if total else 0.0
    positive_rows = [row for row in results if not row["expected_no_match"]]
    negative_rows = [row for row in results if row["expected_no_match"]]
    recall_at_k = sum(1 for row in positive_rows if row["passed"]) / len(positive_rows) if positive_rows else 0.0
    mrr = sum(float(row.get("reciprocal_rank") or 0) for row in positive_rows) / len(positive_rows) if positive_rows else 0.0
    no_match_accuracy = sum(1 for row in negative_rows if row["passed"]) / len(negative_rows) if negative_rows else 0.0
    false_positive_rate = sum(1 for row in negative_rows if not row["passed"]) / len(negative_rows) if negative_rows else 0.0
    config_snapshot = {
        "embedding_model": get_kb(conn, kb_id).get("embedding_model"),
        "rerank_model": get_kb(conn, kb_id).get("rerank_model"),
        "effective_min_threshold": max(0.35, float(get_kb(conn, kb_id).get("score_threshold") or 0.35)),
    }
    with conn.cursor() as cur:
        cur.execute(
            """INSERT INTO rag_eval_runs
            (kb_id,total_cases,passed_cases,pass_rate,average_duration_ms,recall_at_k,
            mean_reciprocal_rank,no_match_accuracy,false_positive_rate,config_json,results_json,created_by)
            VALUES (%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s)""",
            (kb_id, total, passed, pass_rate, average_ms, recall_at_k, mrr, no_match_accuracy,
             false_positive_rate, json.dumps(config_snapshot, ensure_ascii=False),
             json.dumps(results, ensure_ascii=False), username),
        )
        run_id = int(cur.lastrowid)
    conn.commit()
    return {
        "run_id": run_id,
        "total_cases": total,
        "passed_cases": passed,
        "pass_rate": round(pass_rate, 6),
        "average_duration_ms": average_ms,
        "positive_cases": len(positive_rows),
        "negative_cases": len(negative_rows),
        "recall_at_k": round(recall_at_k, 6),
        "mean_reciprocal_rank": round(mrr, 6),
        "no_match_accuracy": round(no_match_accuracy, 6),
        "false_positive_rate": round(false_positive_rate, 6),
        "config": config_snapshot,
        "items": results,
    }


def list_eval_runs(conn: Any, kb_id: int, limit: int = 20) -> List[Dict[str, Any]]:
    with conn.cursor() as cur:
        cur.execute(
            """SELECT id,kb_id,total_cases,passed_cases,pass_rate,average_duration_ms,recall_at_k,
            mean_reciprocal_rank,no_match_accuracy,false_positive_rate,config_json,
            results_json,created_by,created_at FROM rag_eval_runs WHERE kb_id=%s ORDER BY id DESC LIMIT %s""",
            (kb_id, min(max(1, limit), 100)),
        )
        rows = list(cur.fetchall())
    for row in rows:
        try:
            row["items"] = json.loads(row.pop("results_json") or "[]")
        except (TypeError, ValueError):
            row["items"] = []
        try:
            row["config"] = json.loads(row.pop("config_json") or "{}")
        except (TypeError, ValueError):
            row["config"] = {}
    return rows


def migrate_legacy_sqlite(conn: Any, sqlite_path: Path) -> Dict[str, int]:
    """Import the old FTS rows as one source document without deleting it."""
    import sqlite3

    kb_id = ensure_default_kb(conn)
    with conn.cursor() as cur:
        cur.execute("SELECT COUNT(*) AS c FROM rag_documents WHERE source_type='legacy'")
        if int((cur.fetchone() or {}).get("c") or 0) > 0 or not sqlite_path.exists():
            return {"kb_id": kb_id, "documents": 0, "chunks": 0}
    legacy = sqlite3.connect(str(sqlite_path))
    legacy.row_factory = sqlite3.Row
    try:
        rows = [dict(row) for row in legacy.execute("SELECT * FROM rag_docs").fetchall()]
    finally:
        legacy.close()
    if not rows:
        return {"kb_id": kb_id, "documents": 0, "chunks": 0}
    with conn.cursor() as cur:
        cur.execute(
            """INSERT INTO rag_documents
            (kb_id,name,source_type,source_uri,status,progress,chunk_count,char_count,created_by)
            VALUES (%s,%s,'legacy','legacy-sqlite','pending',0,%s,%s,'migration')""",
            (kb_id, "原 RAG 种子知识迁移", len(rows), sum(len(str(row.get("content") or "")) for row in rows)),
        )
        document_id = int(cur.lastrowid)
        for index, row in enumerate(rows):
            content = "\n".join(
                part
                for part in [
                    f"攻击类型：{row.get('attack_type') or '通用'}",
                    str(row.get("content") or ""),
                    f"判定依据：{row.get('evidence') or ''}",
                    f"处置建议：{row.get('mitigation') or ''}",
                    f"来源：{row.get('source') or ''}",
                ]
                if part.strip("：\n ")
            )
            cur.execute(
                """INSERT INTO rag_chunks
                (kb_id,document_id,chunk_index,title_path,content,token_count,vector_key)
                VALUES (%s,%s,%s,%s,%s,%s,%s)""",
                (kb_id, document_id, index, str(row.get("title") or "历史知识")[:600], content, max(1, len(content) // 2), uuid.uuid4().hex),
            )
    conn.commit()
    return {"kb_id": kb_id, "documents": 1, "chunks": len(rows), "document_id": document_id}


def index_pending_document(conn: Any, data_dir: Path, api_config: Dict[str, Any], document_id: int) -> Dict[str, Any]:
    """Vectorize chunks created by the compatibility migration."""
    with conn.cursor() as cur:
        cur.execute("SELECT * FROM rag_documents WHERE id=%s", (document_id,))
        document = cur.fetchone()
        cur.execute("SELECT * FROM rag_chunks WHERE document_id=%s ORDER BY chunk_index", (document_id,))
        chunks = list(cur.fetchall())
    if not document or not chunks:
        raise KeyError("待索引文档不存在")
    vectors = []
    for offset in range(0, len(chunks), 10):
        vectors.extend(embed_texts([row["content"] for row in chunks[offset : offset + 10]], api_config))
    table = _lance_table(data_dir, int(api_config["embedding_dimensions"]), create=True)
    table.add(
        [
            {
                "vector_key": str(row["vector_key"]),
                "kb_id": int(row["kb_id"]),
                "document_id": int(row["document_id"]),
                "chunk_id": int(row["id"]),
                "title": str(row.get("title_path") or document["name"]),
                "content": str(row["content"]),
                "vector": vector,
            }
            for row, vector in zip(chunks, vectors)
        ]
    )
    with conn.cursor() as cur:
        cur.execute("UPDATE rag_documents SET status='ready',progress=100,error_message=NULL WHERE id=%s", (document_id,))
    conn.commit()
    return {"id": document_id, "chunk_count": len(chunks), "status": "ready"}
